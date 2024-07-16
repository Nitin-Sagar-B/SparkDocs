import streamlit as st
import os
import fitz  # PyMuPDF
from pptx import Presentation  # For handling PPTX files
import docx  # For handling DOCX files
from langchain_core.messages import HumanMessage
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.text_splitter import RecursiveCharacterTextSplitter
import re

# Define the correct passkey
correct_passkey = "godsparky1237"  # Replace with your actual passkey

# Function to read a PDF file and extract its text content
def read_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text

# Function to read a PPTX file and extract its text content
def read_pptx(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to read a DOCX file and extract its text content
def read_docx(file):
    doc = docx.Document(file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

# Function to split text into smaller chunks
def split_text(text, chunk_size=5000000, chunk_overlap=100):
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return splitter.split_text(text)

# Function to ask a question and get a response
def ask_question(question, text_chunks, google_api_key):
    # Initialize the model with the provided API key
    model = ChatGoogleGenerativeAI(model="gemini-1.5-flash-latest", google_api_key=google_api_key)
    
    # Clean question for keywords extraction
    cleaned_question = re.sub(r'[^\w\s]', '', question.lower())
    keywords = cleaned_question.split()
    
    # Filter relevant text chunks based on keywords
    relevant_chunks = []
    for chunk in text_chunks:
        chunk_lower = chunk.lower()
        if any(keyword in chunk_lower for keyword in keywords):
            relevant_chunks.append(chunk)
    
    # If no relevant chunks found, use all chunks
    if not relevant_chunks:
        relevant_chunks = text_chunks
    
    # Prepare prompts and get responses
    prompt_template = (
        "system: You are an outstanding AI tutor known for your expertise and clarity. Your role is to assist students by providing clear, concise, and comprehensive answers to their assignment questions using the provided content from the uploaded document. Use your extensive knowledge only if necessary.\n"
        "user: Context: {context}\n\n"
        "Question: {question}\n\n"
        "assistant: \n"
        "Your answers should be:\n"
        "- Well-structured\n"
        "- Very detailed\n"
        "- Comprehensive for the topic\n"
        "- contain short and easy code snippets if needed in context to the topic of the document"
        "- Include relevant examples\n"
        "- Easy to understand for everyone\n"
        "- Free of any repetitive content and unnecessary length\n"
        "- Appropriate for an undergrad engineering AI and Data Science student\n"
        "Begin the answer by writing the question and then the answer"
        "Provide a clear, concise and comprehensive answer based on the PDF context"
    )
    
    full_response = []

    
    for chunk in relevant_chunks:
        prompt = prompt_template.format(context=chunk, question=question)
            
        message = HumanMessage(content=prompt)
        response = model.stream([message])
        
        response_texts = [chunk.content for chunk in response]
        full_response.extend(response_texts)
    
    return ''.join(full_response)

st.set_page_config(layout="wide", page_title="SparkDocs", page_icon=":robot_face:")  # Set app to wide mode and customize title and icon

# Load custom CSS for styling
with open("style.css") as css:
    st.markdown(f'<style>{css.read()}</style>', unsafe_allow_html=True)

# Authentication section
if 'authenticated' not in st.session_state:
    st.session_state.authenticated = False

if not st.session_state.authenticated:
    st.title("Welcome to SparkDocs")

    # Add an input field for the user to enter the passkey
    passkey = st.text_input("🚀 Ready to dive into document analysis? Authenticate with the passkey (hint: it's definitely not 'password123' 😜)", type="password")
    
    # Check if the passkey matches the correct_passkey
    if st.button("Submit"):
        if passkey == correct_passkey:
            st.session_state.authenticated = True
            st.success("Authentication successful! Welcome aboard! 🎉")
            st.rerun()  # Refresh the app to show the main content
        else:
            st.error("Authentication failed. Please enter a valid passkey. 🛑")
    
    st.image("cool logo.png", width=600, caption="Developed by Sparky Labs")

else:
    # Input for API key
    api_key = st.text_input("Enter your Google API key", type="password")

    if not api_key:
        st.warning("Please enter your API key to continue.")
    else:
        # Now you can proceed with your app content here
        st.title("SparkDocs")
        st.subheader("AI-Powered Document Q&A System")
        st.write("Upload a PDF, PPTX, or DOCX file, ask questions, and get answers in context with the document content.")

        uploaded_file = st.file_uploader("Upload Document", type=["pdf", "pptx", "docx"])

        if uploaded_file is not None:
            file_type = uploaded_file.type
            if file_type == "application/pdf":
                doc_text = read_pdf(uploaded_file)
            elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                doc_text = read_pptx(uploaded_file)
            elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc_text = read_docx(uploaded_file)
            else:
                st.error("Unsupported file type!")
                doc_text = ""
            
            if doc_text:
                text_chunks = split_text(doc_text)
                
                question = st.text_input("Enter your question:")
                
                if st.button("Ask"):
                    with st.spinner("Processing..."):
                        answer = ask_question(question, text_chunks, api_key)
                        st.markdown(answer)
